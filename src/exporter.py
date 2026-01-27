from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import os

class NewsExporter:
    def __init__(self):
        pass

    def _clean_text(self, text):
        """
        Sanitizes text to be compatible with FPDF's latin-1 encoding.
        Replaces problematic characters.
        """
        if not text:
            return ""
        # Replace common unicode chars
        replacements = {
            '\u2018': "'", '\u2019': "'", '\u201c': '"', '\u201d': '"',
            '\u2013': '-', '\u2014': '-', '\u2026': '...',
            '\u00f1': 'n', '\u00d1': 'N'
        }
        for k, v in replacements.items():
            text = text.replace(k, v)

        # Fallback: encode to latin-1, replace errors with ?
        return text.encode('latin-1', 'replace').decode('latin-1')

    def to_pdf(self, articles, filename="report.pdf"):
        """
        Generates a PDF report from the list of articles.
        """
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()

        # Title Page
        pdf.set_font("Arial", 'B', 24)
        pdf.cell(0, 20, "Daily News Report", ln=True, align='C')
        pdf.set_font("Arial", size=12)
        pdf.cell(0, 10, "Generated for Journalists", ln=True, align='C')
        pdf.ln(20)

        # Articles
        for i, article in enumerate(articles, 1):
            if not article.get('title'):
                continue

            title = self._clean_text(article['title'])
            summary = self._clean_text(article.get('summary', 'No summary available.'))
            authors = self._clean_text(", ".join(article.get('authors', [])))
            date_str = str(article.get('publish_date')) if article.get('publish_date') else "Date Unknown"
            url = self._clean_text(article['url'])

            # Article Title
            pdf.set_font("Arial", 'B', 16)
            pdf.multi_cell(0, 10, f"{i}. {title}")

            # Metadata
            pdf.set_font("Arial", 'I', 10)
            meta_text = f"Published: {date_str} | Authors: {authors}"
            pdf.cell(0, 8, meta_text, ln=True)

            # Link
            pdf.set_text_color(0, 0, 255)
            # FPDF's link support is tricky with multi_cell, but simple cell works
            # We just print the URL text and make it clickable if supported by the viewer,
            # or use the `link` argument in cell()
            try:
                 pdf.cell(0, 8, "Read Original Article", ln=True, link=url)
            except:
                 pdf.cell(0, 8, f"Source: {url}", ln=True)

            pdf.set_text_color(0, 0, 0)

            pdf.ln(5)

            # Summary
            pdf.set_font("Arial", size=12)
            summary = summary.replace('\n', ' ')
            pdf.multi_cell(0, 8, summary)

            pdf.ln(15)

        pdf.output(filename)
        return filename

    def to_pptx(self, articles, filename="presentation.pptx"):
        """
        Generates a PowerPoint presentation from the list of articles.
        """
        prs = Presentation()

        # Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Daily News Digest"
        subtitle.text = "Generated Report"

        # Content Slides
        bullet_slide_layout = prs.slide_layouts[1]

        for article in articles:
            if not article.get('title'):
                continue

            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            # Use raw text for PPTX (it handles unicode better)
            title_shape.text = article['title']

            tf = body_shape.text_frame
            summary = article.get('summary', 'No summary available.')
            if len(summary) > 600:
                summary = summary[:600] + "..."
            tf.text = summary

            # Add link
            p = tf.add_paragraph()
            p.text = "Source: " + article['url']
            p.font.size = Pt(10)
            p.level = 0

        prs.save(filename)
        return filename

if __name__ == "__main__":
    # Test data
    test_articles = [
        {
            'title': 'Test Article with ñ and utf8: — “quote”',
            'authors': ['Jules'],
            'publish_date': '2023-10-27',
            'summary': 'This is a summary of the first test article. It talks about important things happening in the Philippines.',
            'url': 'http://example.com/1'
        },
        {
            'title': 'Test Article 2',
            'authors': [],
            'publish_date': None,
            'summary': 'Another summary here. Very concise.',
            'url': 'http://example.com/2'
        }
    ]

    exporter = NewsExporter()
    print("Generating PDF...")
    try:
        pdf_path = exporter.to_pdf(test_articles, "test_report.pdf")
        print(f"Saved to {pdf_path}")
    except Exception as e:
        print(f"PDF Error: {e}")

    print("Generating PPTX...")
    try:
        pptx_path = exporter.to_pptx(test_articles, "test_presentation.pptx")
        print(f"Saved to {pptx_path}")
    except Exception as e:
        print(f"PPTX Error: {e}")
